#!/usr/bin/env python3
"""
Style Tableau-exported PowerPoint decks with your brand template.

Requires:
  pip install python-pptx Pillow easyocr

Usage:
  python style_tableau_pptx.py --input tableau_export.pptx --template brand_template.pptx --output styled.pptx \
      --logo /path/logo.png --title-case smart --fit contain --max-margin 32 --footer "The Navigators — Confidential"
      
Note: First run will download OCR language models (~100MB).
"""

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw
import os
from typing import Optional
import easyocr
import numpy as np

def _apply_box_shadow_xml(picture_shape, transparency=0.8, blur_pt=15, angle_deg=34, distance_pt=3, color=(0, 0, 0)):
    """
    DrawingML fallback: inject a:outerShdw into p:spPr/a:effectLst for a picture.
    Uses EMUs for blur/distance and 1/60000 deg for direction per spec.
    """
    try:
        pic_elm = picture_shape._element  # CT_Picture
        # Ensure there is a p:spPr element
        spPr = pic_elm.find(qn('p:spPr'))
        if spPr is None:
            spPr = OxmlElement('p:spPr')
            pic_elm.append(spPr)
        # Ensure there is an a:effectLst
        effectLst = spPr.find(qn('a:effectLst'))
        if effectLst is None:
            effectLst = OxmlElement('a:effectLst')
            spPr.append(effectLst)
        # Remove existing outerShdw to avoid duplicates
        for child in list(effectLst):
            if child.tag == qn('a:outerShdw'):
                effectLst.remove(child)

        outer = OxmlElement('a:outerShdw')
        outer.set('blurRad', str(int(Pt(blur_pt))))
        outer.set('dist', str(int(Pt(distance_pt))))
        outer.set('dir', str(int(angle_deg * 60000)))
        outer.set('algn', 'ctr')
        # Color with alpha (alpha is opacity; transparency 0.8 => alpha 0.2)
        rgb_hex = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
        srgb = OxmlElement('a:srgbClr')
        srgb.set('val', rgb_hex)
        alpha = OxmlElement('a:alpha')
        alpha.set('val', str(int((1.0 - float(transparency)) * 100000)))
        srgb.append(alpha)
        outer.append(srgb)
        effectLst.append(outer)
        return True
    except Exception as e:
        print(f"Warning: XML shadow fallback failed: {e}")
        return False

def apply_box_shadow(shape, transparency=0.8, blur_pt=15, angle_deg=34, distance_pt=3, color=(0, 0, 0)):
    """
    Apply a drop shadow to a shape (picture) using PowerPoint's native shadow.
    Parameters mirror the UI: transparency (0..1), blur/distance in points, angle in degrees.
    """
    try:
        sh = shape.shadow
        # Ensure we do not inherit from theme so settings take effect
        sh.inherit = False
        try:
            sh.visible = True
        except Exception:
            pass
        # PowerPoint expects 0..1 transparency in python-pptx
        sh.transparency = float(transparency)
        sh.blur = Pt(blur_pt)
        sh.angle = int(angle_deg)
        sh.distance = Pt(distance_pt)
        try:
            # Some versions expose fore_color; fallback to color
            sh.fore_color.rgb = RGBColor(*color)
        except Exception:
            try:
                sh.color.rgb = RGBColor(*color)
            except Exception:
                pass
        # Also apply XML fallback for reliability on pictures
        _apply_box_shadow_xml(shape, transparency, blur_pt, angle_deg, distance_pt, color)
    except Exception as e:
        # Best-effort: if the template or library doesn't support shadow here, use XML fallback
        print(f"Warning: could not apply shadow via API: {e}")
        _apply_box_shadow_xml(shape, transparency, blur_pt, angle_deg, distance_pt, color)

def add_rounded_corners(image_path, radius_px=5):
    """
    Add rounded corners to an image and save it.
    Modifies the image file in place.
    """
    try:
        img = Image.open(image_path).convert("RGBA")
        
        # Create a mask for rounded corners
        mask = Image.new('L', img.size, 0)
        draw = ImageDraw.Draw(mask)
        draw.rounded_rectangle([(0, 0), img.size], radius=radius_px, fill=255)
        
        # Apply the mask
        output = Image.new('RGBA', img.size)
        output.paste(img, (0, 0))
        output.putalpha(mask)
        
        # Save back to the same path
        output.save(image_path, 'PNG')
    except Exception as e:
        print(f"Failed to add rounded corners: {e}")

def add_footer(slide, prs, text, color=None, size=10):
    if not text:
        return
    left = Inches(0.7)
    width = Inches(9)
    top = prs.slide_height - Inches(0.5)
    height = Inches(0.3)
    tx = slide.shapes.add_textbox(left, top, width, height)
    p = tx.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    run = p.runs[0]
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = RGBColor(*color)

def smart_titlecase(s: str) -> str:
    """Basic smart title-case with small-word exceptions."""
    if not s:
        return s
    small = set("and or the a an in on for of at to vs via with from".split())
    words = s.split()
    def fix(i,w):
        lw = w.lower()
        if i not in (0, len(words)-1) and lw in small:
            return lw
        return w.capitalize()
    return " ".join(fix(i,w) for i,w in enumerate(words))

def apply_title_case(text: str, mode: str) -> str:
    """Apply requested title casing mode: smart, camel, upper, lower."""
    if not text:
        return text
    mode_lc = (mode or "").lower()
    if mode_lc == "smart":
        return smart_titlecase(text)
    if mode_lc == "camel":
        # naive title case of each word, keeping spaces
        return " ".join(w.capitalize() for w in text.split())
    if mode_lc == "upper":
        return text.upper()
    if mode_lc == "lower":
        return text.lower()
    return text

def guess_title_from_slide(slide) -> Optional[str]:
    # Try to read the first text shape as a title
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                return txt.splitlines()[0][:120]
    return None

def extract_title_from_image(image_path, reader) -> Optional[str]:
    """
    Use OCR to extract the title from the top-left corner of an image.
    Assumes the title is in the top-left 40% width x 15% height area.
    """
    try:
        with Image.open(image_path) as img:
            width, height = img.size
            # Crop to top-left corner where Tableau titles typically are
            # Left 40% of width, top 15% of height
            top_left = img.crop((0, 0, int(width * 0.4), int(height * 0.05)))
            
            # Convert PIL Image to numpy array for easyocr
            img_array = np.array(top_left)
            
            # Run OCR on the cropped portion (without paragraph mode)
            result = reader.readtext(img_array)
            
            if result:
                # easyocr returns list of (bbox, text, confidence)
                texts = [text for (bbox, text, conf) in result if conf > 0.3]
                if texts:
                    # Join all text pieces (in case title is split)
                    # Sort by vertical position (top to bottom)
                    sorted_results = sorted(result, key=lambda x: x[0][0][1])  # sort by y-coordinate
                    title = " ".join([text for (bbox, text, conf) in sorted_results if conf > 0.3])
                    return title.strip()[:120] if title.strip() else None
    except Exception as e:
        print(f"OCR failed: {e}")
    return None

def find_layout(tpl, preferred=("Title Only","Title and Content","Blank")):
    name_to_layout = {l.name: l for l in tpl.slide_layouts}
    for name in preferred:
        if name in name_to_layout:
            return name_to_layout[name]
    # fallback to first
    return tpl.slide_layouts[0]

def add_slide_number(prs):
    # Ensure slide numbers are visible (Office respects this flag)
    try:
        prs.core_properties.revision = str(int(prs.core_properties.revision or "0") + 1)
        el = prs.part.presentation.sldMasterIdLst  # poke to ensure presence
        # We rely on the template having slide numbers enabled in the master.
    except Exception:
        pass

def fit_image_on_blank(slide, prs, pic, left_in=2.5, top_in=1.7, height_in=4.9):
    """
    Reposition + scale the picture and place at provided left/top (inches).
    Sets height to height_in while preserving aspect ratio.
    """
    shp = pic
    # Preserve aspect ratio at requested height
    original_width = shp.width
    original_height = shp.height
    aspect_ratio = original_width / original_height
    
    shp.height = Inches(height_in)
    shp.width = int(shp.height * aspect_ratio)
    
    # Set explicit placement
    shp.left = Inches(left_in)
    shp.top = Inches(top_in)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--input", "-i", required=True, help="Tableau-exported PPTX")
    ap.add_argument("--template", "-t", required=True, help="Brand template PPTX (with masters, theme, slide numbers, etc.)")
    ap.add_argument("--output", "-o", default="styled_output.pptx")
    ap.add_argument("--title-case", choices=["smart","camel","upper","lower"], default="smart")
    ap.add_argument("--title-font-size", type=int, default=28)
    ap.add_argument("--border-radius", type=int, default=10, help="Border radius in pixels for rounded corners on images")
    # Image placement (inches)
    ap.add_argument("--image-left", type=float, default=2.5, help="Left position in inches for the image")
    ap.add_argument("--image-top", type=float, default=1.7, help="Top position in inches for the image")
    # Shadow controls (on by default)
    ap.add_argument("--shadow", dest="shadow", action="store_true", default=True, help="Apply drop shadow to images (default on)")
    ap.add_argument("--no-shadow", dest="shadow", action="store_false", help="Disable drop shadow on images")
    ap.add_argument("--shadow-color", default="000000", help="Shadow RGB hex like 000000")
    ap.add_argument("--shadow-transparency", type=float, default=0.8, help="Shadow transparency 0..1 (default 0.8)")
    ap.add_argument("--shadow-blur", type=int, default=15, help="Shadow blur in points (default 15)")
    ap.add_argument("--shadow-angle", type=int, default=34, help="Shadow angle in degrees (default 34)")
    ap.add_argument("--shadow-distance", type=int, default=3, help="Shadow distance in points (default 3)")
    args = ap.parse_args()

    src = Presentation(args.input)
    tpl = Presentation(args.template)
    out = Presentation(args.template)  # base on template to inherit theme/master
    layout = find_layout(tpl)
    add_slide_number(out)

    # Initialize OCR reader once (first run will download model)
    print("Initializing OCR reader (first run may download language models)...")
    reader = easyocr.Reader(['en'], gpu=False)
    print("OCR reader ready.")

    # Parse shadow color
    shadow_color = None
    try:
        shadow_color = tuple(int(args.shadow_color[i:i+2], 16) for i in (0,2,4))
    except Exception:
        shadow_color = (0, 0, 0)

    for idx, s in enumerate(src.slides, start=1):
        # create new slide using chosen layout
        slide = out.slides.add_slide(layout)

        # 1) Extract image first (needed for OCR)
        image_path = None
        for shp in s.shapes:
            if getattr(shp, "image", None) is not None:
                ext = shp.image.ext or "png"
                tmp = f"_tmp_slide_{idx}.{ext}"
                with open(tmp, "wb") as f:
                    f.write(shp.image.blob)
                image_path = tmp
                break

        # 2) Extract title using OCR from image, or fallback to slide text
        title_text = None
        if image_path and os.path.exists(image_path):
            print(f"Extracting title from image {idx}...")
            title_text = extract_title_from_image(image_path, reader)
        
        # Fallback to text extraction from slide if OCR didn't work
        if not title_text:
            title_text = guess_title_from_slide(s) or f"Dashboard {idx}"
        
        # Limit to a reasonable length to avoid overflows
        title_text = (title_text or "")[:120]
        title_text = apply_title_case(title_text, args.title_case)
        
        # Put title into slide
        if slide.shapes.title:
            slide.shapes.title.text = title_text
            slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(args.title_font_size)
        else:
            tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), out.slide_width - Inches(1.4), Inches(0.6))
            p = tb.text_frame.paragraphs[0]
            p.text = title_text
            p.runs[0].font.size = Pt(args.title_font_size)

        # 3) Add image to slide with rounded corners
        if image_path and os.path.exists(image_path):
            # Add rounded corners to the image
            add_rounded_corners(image_path, radius_px=args.border_radius)
            
            pic = slide.shapes.add_picture(image_path, 0, 0)
            fit_image_on_blank(slide, out, pic, left_in=args.image_left, top_in=args.image_top)
            # Apply drop shadow if enabled
            if args.shadow:
                apply_box_shadow(
                    pic,
                    transparency=args.shadow_transparency,
                    blur_pt=args.shadow_blur,
                    angle_deg=args.shadow_angle,
                    distance_pt=args.shadow_distance,
                    color=shadow_color,
                )
            try:
                os.remove(image_path)
            except Exception:
                pass

    out.save(args.output)
    print(f"✅ Wrote {args.output}")

if __name__ == "__main__":
    main()
